from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN = RGBColor(0x33, 0x69, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)
DARK = RGBColor(0x44, 0x44, 0x44)
MUTED = RGBColor(0x77, 0x77, 0x77)
FAINT = RGBColor(0xAA, 0xAA, 0xAA)
BG = RGBColor(0xFA, 0xFA, 0xF8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xDF, 0xD9)
WARNING = RGBColor(0xC7, 0x75, 0x00)
ERROR = RGBColor(0xC6, 0x28, 0x28)
SUCCESS = RGBColor(0x1B, 0x7A, 0x3D)
GREEN_LIGHT = RGBColor(0x55, 0x8B, 0x2F)

CHARTS = r"C:\Users\User\Desktop\project\llm_from_scratch-main\LOADFORECASINGPROJECT\webapp\static\charts"

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, italic=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs_list, font_name="Calibri", alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(runs_list):
        if isinstance(item, dict):
            text = item.get("text", "")
            opts = item.get("options", {})
        else:
            text = item
            opts = {}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(opts.get("size", 15))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK)
        p.font.name = opts.get("font", font_name)
        p.alignment = alignment
        p.space_after = Pt(opts.get("space_after", 4))
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = Pt(opts.get("size", 15) * line_spacing)
    return txBox

def add_image(slide, path, left, top, width, height=None):
    if height is None:
        height = width * 0.5625
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))

def add_accent_bar(slide, top=0.3, width=0.06):
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(width), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

def add_slide_header(slide, num, context):
    add_accent_bar(slide, 0.35, 0.05)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}  /  {context}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = "Calibri"

def add_slide_title(slide, title, top=0.8):
    add_textbox(slide, 0.5, top, 12.3, 0.8, title, font_size=32, bold=True, color=BLACK)

def add_footer(slide, text):
    add_textbox(slide, 0.5, 7.0, 12, 0.4, text, font_size=10, color=MUTED, font_name="Calibri")

# ======== SLIDE 0: COVER ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x11, 0x11, 0x11))

add_textbox(slide, 0.8, 0.5, 8, 0.5, "KNUST  ×  GridCo", font_size=22, bold=True, color=GREEN_LIGHT)

add_textbox(slide, 0.8, 2.0, 4, 0.35, "FROM ACADEMIC RESEARCH TO PRODUCTION IMPACT", font_size=11, bold=True, color=GREEN_LIGHT)
add_textbox(slide, 0.8, 2.4, 11, 1.6, "DLinear+H10: Load Forecasting\nfor Ghana's National Grid.", font_size=48, bold=True, color=WHITE)
add_textbox(slide, 0.8, 4.2, 9, 0.8, "A joint project between KNUST and GridCo — bringing accurate, data-driven forecasting\nto the national dispatch centre. Saving fuel. Saving money. Keeping the lights on.", font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

add_textbox(slide, 0.8, 6.0, 12, 0.6, "70,228 hours trained    91 MW baseline    6-fold CV    438,000 MWh/year saved", font_size=15, bold=True, color=RGBColor(0xCC, 0xCC, 0xCC))

right_box = slide.shapes.add_shape(1, Inches(9.5), Inches(2.0), Inches(3.2), Inches(3.0))
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
right_box.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
right_box.line.width = Pt(1)

items = [
    ("Project Partners", "KNUST × GridCo"),
    ("Engine", "DLinear + H10 Corrector"),
    ("Operational Savings", "438,000 MWh/year refined"),
    ("Data Span", "Jan 2018 – May 2026"),
]
y = 2.3
for label, val in items:
    add_textbox(slide, 9.8, y, 2.6, 0.25, label, font_size=9, bold=True, color=FAINT)
    add_textbox(slide, 9.8, y + 0.25, 2.6, 0.3, val, font_size=13, bold=True, color=WHITE)
    y += 0.65

# ======== SLIDE 1: THE DATA ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "01", "The Data  —  70k hours from Ghana's national grid")
add_slide_title(slide, "What the Data Tells Us About Ghana's Growing Grid.")

add_image(slide, os.path.join(CHARTS, "demand_distribution.png"), 0.5, 1.9, 6.5, 3.8)

box = slide.shapes.add_shape(1, Inches(7.3), Inches(1.9), Inches(5.5), Inches(2.5))
box.fill.solid()
box.fill.fore_color.rgb = WHITE
box.line.color.rgb = BORDER
box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(7.3), Inches(1.9), Inches(0.06), Inches(2.5))
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = GREEN
left_bar.line.fill.background()

add_rich_textbox(slide, 7.6, 2.0, 5.0, 2.3, [
    {"text": "Here's What We Found", "options": {"size": 20, "bold": True, "color": BLACK, "space_after": 10}},
    {"text": "Demand grew 2.5× in 8 years. From ~1,467 MW to ~3,750 MW. That's 12.4% annual growth. Compare this to Europe: 0–1%. A model trained on 2021 data simply won't work on 2026 data.", "options": {"size": 13, "color": DARK, "space_after": 8}},
    {"text": "Accuracy has a price tag. Every MW of forecast error costs GridCo money — either fuel burned for over-generation that nobody uses, or expensive peaker plants scrambled last-minute for under-generation. A better forecast means leaner dispatch, less waste, lower costs.", "options": {"size": 13, "color": DARK, "space_after": 8}},
    {"text": "What we don't know: electrification (67% → 83% rural access), GDP, urbanisation — these drive demand but we don't track them as features.", "options": {"size": 13, "color": DARK, "space_after": 0}},
])

box2 = slide.shapes.add_shape(1, Inches(7.3), Inches(4.6), Inches(5.5), Inches(1.7))
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = BORDER
box2.line.width = Pt(1)
add_rich_textbox(slide, 7.6, 4.7, 5.0, 1.5, [
    {"text": "Patterns That Jumped Out", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 8}},
    {"text": "→  Daily rhythm: Morning ramp. Evening peak. Overnight trough.", "options": {"size": 12, "color": DARK, "space_after": 2}},
    {"text": "→  Weekly cycle: Weekdays high. Weekends low.", "options": {"size": 12, "color": DARK, "space_after": 2}},
    {"text": "→  Seasons: Dry season ≠ wet season.", "options": {"size": 12, "color": DARK, "space_after": 2}},
    {"text": "→  Growth: Upward drift across all 8 years.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

add_textbox(slide, 0.5, 5.9, 12, 0.3, "Features: demand_mw  ·  temperature_c  ·  hour  ·  date  ·  3.95% filtered for sensor outages", font_size=11, color=MUTED)
add_footer(slide, "Source: ECG SCADA  ·  12.4% CAGR = Compound Annual Growth Rate 2018–2026")

# ======== SLIDE 2: AIMS ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "02", "Aims  —  What we set out to achieve")
add_slide_title(slide, "Five Aims for This Academia–Industry Collaboration.")

# Left column — aims 1 & 2
box = slide.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.3))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(0.06), Inches(2.3))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 0.7, 2.1, 5.5, 2.1, [
    {"text": "1. Accurate, Robust Forecasting", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "Develop a data-driven load forecasting model for Ghana's national grid that handles rapid demand growth (12.4% CAGR), distribution shifts, and limited feature availability — outperforming both operator heuristics and imported deep learning architectures.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

box = slide.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(5.9), Inches(1.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(0.06), Inches(1.8))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 0.7, 4.6, 5.5, 1.6, [
    {"text": "2. Error Correction", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "Go beyond raw predictions. Study residual errors to build second-stage correctors that extract an additional 20–40% accuracy gain from autocorrelated error patterns.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

# Right column — aims 3, 4, 5
box = slide.shapes.add_shape(1, Inches(6.7), Inches(2.0), Inches(6.1), Inches(1.6))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(6.7), Inches(2.0), Inches(0.06), Inches(1.6))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 6.9, 2.1, 5.7, 1.4, [
    {"text": "3. Production Deployment", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "Package research into a web application deployed at GridCo's dispatch centre — replacing a 45-minute Excel workflow with a 1-minute automated process. Keep operators in control.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

box = slide.shapes.add_shape(1, Inches(6.7), Inches(3.8), Inches(6.1), Inches(1.4))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(6.7), Inches(3.8), Inches(0.06), Inches(1.4))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 6.9, 3.9, 5.7, 1.2, [
    {"text": "4. Sustainable AI", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "Prove that accurate grid forecasting does not require GPU infrastructure or massive energy budgets. A 36K-parameter CPU-trained model can beat Transformers on both accuracy and efficiency.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

box = slide.shapes.add_shape(1, Inches(6.7), Inches(5.4), Inches(6.1), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(6.7), Inches(5.4), Inches(0.06), Inches(1.2))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 6.9, 5.5, 5.7, 1.0, [
    {"text": "5. Insights for West African Grids", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "Draw generalisable lessons about load forecasting in rapidly developing electricity grids where stationary data assumptions fail.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

add_footer(slide, "KNUST Department of Electrical Engineering  ·  GridCo National Dispatch Centre")

# ======== SLIDE 3: GRIDCO ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "03", "Current Practice  —  The operator's existing workflow")
add_slide_title(slide, "GridCo's Current Approach: Manual Heuristics → Data-Driven.")

# Card 1
box = slide.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(3.8), Inches(2.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 0.7, 2.1, 3.4, 2.6, [
    {"text": "What They Do Today", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 6}},
    {"text": "Dispatchers use Excel. They find a similar historical day by experience, adjust for temperature, populate the schedule manually.", "options": {"size": 12, "color": DARK, "space_after": 6}},
    {"text": "~45 minutes per forecast. No learning from 70k hours of data. Each forecast starts from scratch.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

# Card 2
box = slide.shapes.add_shape(1, Inches(4.6), Inches(2.0), Inches(3.8), Inches(2.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(4.6), Inches(2.0), Inches(0.06), Inches(2.8))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 4.8, 2.1, 3.4, 2.6, [
    {"text": "We Tested Their Approach", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 6}},
    {"text": "Digitised the similar-day method as a KNN model. Result: 141 MW MAE.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "A simple 2-day weighted trend? 113 MW MAE. Beat the heuristic.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "Conclusion: Data-driven methods outperform human heuristics when given enough history.", "options": {"size": 12, "bold": True, "color": DARK, "space_after": 0}},
])

# Card 3
box = slide.shapes.add_shape(1, Inches(8.7), Inches(2.0), Inches(4.1), Inches(2.8))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xEE); box.line.color.rgb = RGBColor(0xC8, 0xDD, 0xB8); box.line.width = Pt(1)
add_rich_textbox(slide, 8.9, 2.1, 3.7, 2.6, [
    {"text": "So What Did We Build?", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 6}},
    {"text": "A system that complements the operator — doesn't replace them.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "Forecast auto-generates in ~1 minute. Fills the dispatch schedule automatically.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "Operator stays in control. The forecast is a suggestion — accept, edit, or override.", "options": {"size": 12, "bold": True, "color": DARK, "space_after": 0}},
])

# Bottom stats
add_textbox(slide, 0.5, 5.3, 2.5, 0.5, "Manual time", font_size=9, bold=True, color=MUTED)
add_textbox(slide, 0.5, 5.7, 2.5, 0.5, "~45 min", font_size=22, bold=True, color=WARNING)
add_textbox(slide, 3.2, 5.3, 2.5, 0.5, "KNN baseline", font_size=9, bold=True, color=MUTED)
add_textbox(slide, 3.2, 5.7, 2.5, 0.5, "141 MW", font_size=22, bold=True, color=FAINT)
add_textbox(slide, 5.9, 5.3, 2.5, 0.5, "Simple trend", font_size=9, bold=True, color=MUTED)
add_textbox(slide, 5.9, 5.7, 2.5, 0.5, "113 MW", font_size=22, bold=True, color=FAINT)
add_textbox(slide, 8.6, 5.3, 3, 0.5, "Our system", font_size=9, bold=True, color=MUTED)
add_textbox(slide, 8.6, 5.7, 3, 0.5, "~1 min  ·  68.9 MW", font_size=22, bold=True, color=SUCCESS)

add_footer(slide, "KNN = k-Nearest Neighbours  ·  MAE = Mean Absolute Error")

# ======== SLIDE 3: INITIAL WORK ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "04", "Foundations  —  Understanding data before modelling")
add_slide_title(slide, "We Started Simple. Here's What We Found First.")

box = slide.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 0.7, 2.1, 5.5, 2.6, [
    {"text": "Patterns Jumped Out", "options": {"size": 20, "bold": True, "color": BLACK, "space_after": 8}},
    {"text": "→  Daily rhythm: Morning ramp. Evening peak. Overnight trough. Every day.", "options": {"size": 13, "color": DARK, "space_after": 3}},
    {"text": "→  Weekly rhythm: Weekdays high. Weekends low. Predictable.", "options": {"size": 13, "color": DARK, "space_after": 3}},
    {"text": "→  Seasons: Dry season ≠ wet season. Load changes with the calendar.", "options": {"size": 13, "color": DARK, "space_after": 3}},
    {"text": "→  Growth: Upward drift across all 8 years. +12.4% CAGR in the raw data.", "options": {"size": 13, "color": DARK, "space_after": 3}},
    {"text": "On temperature: It affects load — but it's slow-moving and already baked into hour and month features.", "options": {"size": 12, "color": MUTED, "space_after": 0}},
])

box = slide.shapes.add_shape(1, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(6.7), Inches(2.0), Inches(0.06), Inches(2.8))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 6.9, 2.1, 5.7, 2.6, [
    {"text": "Resisted the Urge to Go Big", "options": {"size": 20, "bold": True, "color": BLACK, "space_after": 8}},
    {"text": "Most projects reach for deep learning first. We didn't. We asked: how far can simple methods take us?", "options": {"size": 13, "color": DARK, "space_after": 8}},
    {"text": "First Try:  Classical decomposition (trend + seasonality)  →  3–4% MAPE", "options": {"size": 13, "bold": True, "color": GREEN_LIGHT, "space_after": 4}},
    {"text": "Second Try:  Weighted trend (yesterday + last week)  →  113 MW MAE", "options": {"size": 13, "bold": True, "color": GREEN_LIGHT, "space_after": 6}},
    {"text": "Lesson: Simple models already capture most of the signal. Deep learning's value must be proven — not assumed.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

add_textbox(slide, 0.5, 5.3, 12.3, 0.6, "Our rule: Establish a strong simple baseline first. If classical decomposition hits 3–4% MAPE, don't jump to a transformer. Prove the marginal value of complexity.", font_size=13, color=DARK)
add_footer(slide, "MAPE = Mean Absolute Percentage Error  ·  L1 = last 24 hours  ·  L7 = last 7 days")

# ======== SLIDE 4: DL BENCHMARKING ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "05", "Benchmarking  —  6-fold expanding-window CV across 8 years")
add_slide_title(slide, "DLinear Beat Transformers, LSTMs, and CNNs — Here's How.")

# Table
rows = 5
cols = 4
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(6.0), Inches(2.0)).table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(1.2)
table.columns[3].width = Inches(1.1)

headers = ["Model", "Mean MAE", "+TIDE", "Gain"]
data = [
    ["DLinear  ·  36K params", "91 MW", "67 MW", "−26%"],
    ["CNN (WaveNet)  ·  1.2M", "97 MW", "74 MW", "−23%"],
    ["LSTM  ·  840K", "102 MW", "79 MW", "−23%"],
    ["Transformer  ·  2.1M", "109 MW", "82 MW", "−24%"],
]

for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = FAINT
        p.font.name = "Calibri"

for ri, row_data in enumerate(data):
    for ci, val in enumerate(row_data):
        cell = table.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = GREEN if ci < 2 else DARK
            else:
                p.font.color.rgb = DARK

add_image(slide, os.path.join(CHARTS, "fold_mae.png"), 0.5, 4.2, 6.0, 2.6)

box = slide.shapes.add_shape(1, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(6.8), Inches(2.0), Inches(0.06), Inches(4.8))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()

add_rich_textbox(slide, 7.0, 2.2, 5.6, 4.5, [
    {"text": "Why DLinear Won", "options": {"size": 20, "bold": True, "color": BLACK, "space_after": 8}},
    {"text": "DLinear = Decomposition Linear. It splits load into two parts:", "options": {"size": 13, "color": DARK, "space_after": 6}},
    {"text": "→  Trend — gradual direction (is demand going up or down?)", "options": {"size": 13, "color": DARK, "space_after": 3}},
    {"text": "→  Seasonal — daily/weekly rhythm (morning peak, weekend dip)", "options": {"size": 13, "color": DARK, "space_after": 8}},
    {"text": "Each gets a learned weight plus calendar info. They sum to the forecast. No attention. No recurrence. No convolutions.", "options": {"size": 13, "color": DARK, "space_after": 8}},
    {"text": "Why it works: Load patterns aren't that complex. DLinear's 36K parameters capture the signal. The other models (1–2M params) fit noise. DLinear trains on a laptop CPU in minutes. The simplest model won — and every MW DLinear saves vs the Transformer means real fuel not burned at GridCo's dispatch.", "options": {"size": 13, "color": DARK, "space_after": 0}},
])

add_footer(slide, "6-fold expanding window  ·  8 years of data  ·  TIDE = Trend-adjusted Iterative Debiasing Engine")

# ======== SLIDE 5: DLINAR RESULTS ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "06", "Results  —  Raw DLinear before any correction")
add_slide_title(slide, "DLinear's Performance Before Correction.")

# Horizon table
rows = 4
cols = 3
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(4.5), Inches(1.8)).table
table.columns[0].width = Inches(2.3)
table.columns[1].width = Inches(1.1)
table.columns[2].width = Inches(1.1)

h_data = [
    ["Horizon", "MAE", "MAPE"],
    ["Day-Ahead (24h)", "121 MW", "4.22%"],
    ["Week-Ahead (168h)", "163 MW", "5.30%"],
    ["Month-Ahead (720h)", "104 MW", "3.40%"],
]
for ri, row in enumerate(h_data):
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            if ri == 0:
                p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = FAINT
            else:
                p.font.size = Pt(13); p.font.color.rgb = DARK
            p.font.name = "Calibri"

add_textbox(slide, 0.5, 4.0, 4.5, 0.5, "One model handles all three horizons. 720h lower than 168h because longer horizons average out daily variation.", font_size=11, color=MUTED)

# MAE vs MAPE card
box = slide.shapes.add_shape(1, Inches(5.3), Inches(2.0), Inches(4.0), Inches(2.5))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(5.3), Inches(2.0), Inches(0.06), Inches(2.5))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 5.5, 2.1, 3.6, 2.3, [
    {"text": "MAE vs MAPE — Why Both?", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 6}},
    {"text": "MAE: Average error in megawatts. What GridCo dispatchers feel — they dispatch in MW.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "MAPE: Error as a percentage. Useful across different load levels. 121 MW is ~3.9% at peak but larger at low load.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

# 2024 spike card
box = slide.shapes.add_shape(1, Inches(9.6), Inches(2.0), Inches(3.2), Inches(2.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xEE); box.line.color.rgb = RGBColor(0xC8, 0xDD, 0xB8); box.line.width = Pt(1)
add_rich_textbox(slide, 9.8, 2.1, 2.8, 2.3, [
    {"text": "The 2024 Spike", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 6}},
    {"text": "Fold 4 hit 226 MW — more than double the average.", "options": {"size": 12, "color": ERROR, "space_after": 4}},
    {"text": "COVID recovery broke the pattern. A model that worked in 2023 failed in 2024.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "This is why 6-fold CV matters: one number hides regime-dependent failure.", "options": {"size": 12, "bold": True, "color": DARK, "space_after": 0}},
])

# Per-fold table
rows = 6
cols = 4
table = slide.shapes.add_table(rows, cols, Inches(5.3), Inches(4.8), Inches(7.5), Inches(2.0)).table
table.columns[0].width = Inches(0.8)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(1.0)
table.columns[3].width = Inches(3.7)

fold_data = [
    ["Fold", "Test Year", "MAE", "Context"],
    ["F1", "2021", "166 MW", "COVID recovery disruption"],
    ["F2", "2022", "107 MW", "Stable recovery"],
    ["F3", "2023", "111 MW", "Moderate growth"],
    ["F4", "2024", "226 MW", "Growth discontinuity — spike"],
    ["F5", "2025", "92 MW", "Pattern stabilized"],
]
for ri, row in enumerate(fold_data):
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            if ri == 0:
                p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = FAINT
            else:
                p.font.size = Pt(12); p.font.color.rgb = DARK
                if ri == 4 and ci == 2:
                    p.font.color.rgb = ERROR; p.font.bold = True
            p.font.name = "Calibri"

add_footer(slide, "6-fold mean: 110 MW  ·  Range: 92–226 MW")

# ======== SLIDE 6: ERROR CORRECTION ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "07", "Correction  —  Learning from prediction errors")
add_slide_title(slide, "Good Results — But Errors Told a Different Story.")

box = slide.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(6.0), Inches(1.5))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
left_bar = slide.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(0.06), Inches(1.5))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = GREEN; left_bar.line.fill.background()
add_rich_textbox(slide, 0.7, 2.1, 5.6, 1.3, [
    {"text": "Errors Are Not Random", "options": {"size": 18, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "We checked the residuals. Normally a good model's errors look like noise. Ours didn't. Correlation of 0.79 between consecutive hours. If we over-predicted hour 0, we'd over-predict hour 1 too. Errors have structure. That makes them learnable.", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

add_image(slide, os.path.join(CHARTS, "error_correction.png"), 6.8, 2.0, 6.0, 2.6)

# TIDE card
box = slide.shapes.add_shape(1, Inches(0.5), Inches(3.8), Inches(3.8), Inches(1.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 0.7, 3.9, 3.4, 1.6, [
    {"text": "TIDE — A Smart Filter", "options": {"size": 16, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "Zero parameters. No training. Dampens recent errors using a weighted average. Learns how much to trust the latest mistake vs. the long-term trend.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "Result: 95.5 MW  (−20.9%)", "options": {"size": 14, "bold": True, "color": WARNING, "space_after": 0}},
])

# ARD card
box = slide.shapes.add_shape(1, Inches(4.6), Inches(3.8), Inches(3.8), Inches(1.8))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 4.8, 3.9, 3.4, 1.6, [
    {"text": "ARD — Probabilistic Fix", "options": {"size": 16, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "Bayesian model that learns persistent bias: 'DLinear always over-predicts Tuesday 19h by 15 MW.' Only works with true error access (sequential mode).", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "Result: 68.9 MW  (−40.4%)", "options": {"size": 14, "bold": True, "color": GREEN, "space_after": 0}},
])

# Key finding card
box = slide.shapes.add_shape(1, Inches(8.7), Inches(3.8), Inches(4.1), Inches(1.8))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xEE); box.line.color.rgb = RGBColor(0xC8, 0xDD, 0xB8); box.line.width = Pt(1)
add_rich_textbox(slide, 8.9, 3.9, 3.7, 1.6, [
    {"text": "The Critical Finding", "options": {"size": 16, "bold": True, "color": BLACK, "space_after": 4}},
    {"text": "TIDE works because it has access to recent errors. ARD without error history? Useless.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "A streaming model — correcting hour-by-hour as actuals arrive — unlocks the full 40% gain. And it does so with near-zero operational cost while saving GridCo real fuel and money on dispatch.", "options": {"size": 12, "bold": True, "color": DARK, "space_after": 0}},
])

add_textbox(slide, 0.5, 5.9, 12.3, 0.4, "All methods: Baseline 115.6  ·  TIDE 95.5 (−21%)  ·  SMA-7d 106  ·  Kalman 98  ·  ARD batch 115  ·  ARD sequential 68.9 (−40%)", font_size=11, color=MUTED)
add_footer(slide, "TIDE = Trend-adjusted Iterative Debiasing Engine  ·  ARD = Automatic Relevance Determination")

# ======== SLIDE 8: ENERGY EFFICIENCY ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "08", "Energy  —  How better accuracy saves real fuel and emissions")
add_slide_title(slide, "Better Forecasts Don't Just Improve Numbers — They Save Fuel.")

add_image(slide, os.path.join(CHARTS, "energy_efficiency.png"), 0.5, 1.9, 12.3, 4.0)

# Box 1
box = slide.shapes.add_shape(1, Inches(0.5), Inches(6.1), Inches(3.8), Inches(1.4))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xEE); box.line.color.rgb = RGBColor(0xC8, 0xDD, 0xB8); box.line.width = Pt(1)
add_rich_textbox(slide, 0.7, 6.2, 3.4, 1.2, [
    {"text": "The Fuel Story. 50 MW less error → 50 MW dispatched correctly every hour. No wasted fuel from over-generation. No expensive peaker scramble from under-generation. At ~0.4 tCO₂/MWh, that's ~20 tCO₂ saved per hour of avoided over-gen.", "options": {"size": 11, "color": DARK, "space_after": 0}},
])

# Box 2
box = slide.shapes.add_shape(1, Inches(4.6), Inches(6.1), Inches(3.8), Inches(1.4))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 4.8, 6.2, 3.4, 1.2, [
    {"text": "The Money Story. GridCo's peaker plants cost $150–300/MWh. A 50 MW under-forecast mistake can cost thousands per hour in fast-ramping generation. DLinear cuts that error by 35–50%. At 12.4% demand growth, the cost of being wrong compounds yearly.", "options": {"size": 11, "color": DARK, "space_after": 0}},
])

# Box 3
box = slide.shapes.add_shape(1, Inches(8.7), Inches(6.1), Inches(4.1), Inches(1.4))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 8.9, 6.2, 3.7, 1.2, [
    {"text": "The Annual Picture. 50 MW avg × 24h × 365d = 438,000 MWh/year of better dispatch decisions. The model trains for 0.003 kWh per fold. The operational savings — fuel, emissions, peaker avoidance — are millions of times larger.", "options": {"size": 11, "color": DARK, "space_after": 0}},
])

add_footer(slide, "Training energy measured with CodeCarbon  ·  Gas turbine: ~40% efficiency, ~0.4 tCO₂/MWh  ·  Peaker fuel ~$150/MWh")

# ======== SLIDE 9: TAKEAWAYS ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_slide_header(slide, "09", "Takeaways  —  What this means for West African grids")
add_slide_title(slide, "What This Work Says About West African Grids.")

# Industrial impact
box = slide.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.5))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 0.7, 2.1, 5.5, 2.3, [
    {"text": "Industrial Impact", "options": {"size": 20, "bold": True, "color": BLACK, "space_after": 6}},
    {"text": "→  Replaced a ~45-minute manual Excel workflow with a ~1-minute automated process", "options": {"size": 13, "color": DARK, "space_after": 4}},
    {"text": "→  Web application deployed at GridCo — forecast auto-fills the dispatch schedule", "options": {"size": 13, "color": DARK, "space_after": 4}},
    {"text": "→  Operators stay in control: accept, edit, or override any forecast", "options": {"size": 13, "color": DARK, "space_after": 4}},
    {"text": "→  Designed for semi-annual retraining as new data arrives", "options": {"size": 13, "color": DARK, "space_after": 4}},
    {"text": "→  Fuel savings: 50 MW better accuracy = ~438,000 MWh/year smarter dispatch — less fuel burned, lower emissions, cheaper operations", "options": {"size": 13, "color": DARK, "space_after": 0}},
])

# For West African grids
box = slide.shapes.add_shape(1, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xEE); box.line.color.rgb = RGBColor(0xC8, 0xDD, 0xB8); box.line.width = Pt(1)
add_rich_textbox(slide, 6.9, 2.1, 5.7, 2.3, [
    {"text": "For West African Grids", "options": {"size": 20, "bold": True, "color": BLACK, "space_after": 6}},
    {"text": "Stationary data is a myth here. European grids grow 0–1% yearly. Ours: 12.4%. Models from mature grids degrade fast.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "6-fold CV is non-negotiable. Model rankings change year to year (2024's 226 MW spike). A single test year can mislead.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "Accuracy is fuel efficiency. Every MW of forecast error is a MW misdispatched — wasted fuel or expensive reserves. Better forecasts = leaner, cheaper grid.", "options": {"size": 12, "color": DARK, "space_after": 4}},
    {"text": "Streaming is the path forward. The 40% gain proves that error autocorrelation is our strongest signal. Batch can't use it.", "options": {"size": 12, "bold": True, "color": DARK, "space_after": 0}},
])

# Roadmap
box = slide.shapes.add_shape(1, Inches(0.5), Inches(4.9), Inches(3.8), Inches(1.6))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 0.7, 5.0, 3.4, 1.4, [
    {"text": "Short-term", "options": {"size": 14, "bold": True, "color": WARNING, "space_after": 4}},
    {"text": "Deploy robust correctors (TIDE-style) — reliable gains without new infrastructure", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

box = slide.shapes.add_shape(1, Inches(4.6), Inches(4.9), Inches(3.8), Inches(1.6))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 4.8, 5.0, 3.4, 1.4, [
    {"text": "Medium-term", "options": {"size": 14, "bold": True, "color": GREEN_LIGHT, "space_after": 4}},
    {"text": "Incorporate external drivers: electrification rates, GDP, urbanisation as model features", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

box = slide.shapes.add_shape(1, Inches(8.7), Inches(4.9), Inches(4.1), Inches(1.6))
box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = BORDER; box.line.width = Pt(1)
add_rich_textbox(slide, 8.9, 5.0, 3.7, 1.4, [
    {"text": "Long-term", "options": {"size": 14, "bold": True, "color": GREEN, "space_after": 4}},
    {"text": "Streaming DLinear with continuous real-time correction — lightweight model + live error feedback", "options": {"size": 12, "color": DARK, "space_after": 0}},
])

add_textbox(slide, 0.5, 6.7, 10, 0.5, '"The grid will keep changing. Research must change with it. The simplest model won — but streaming correction unlocks what batch cannot touch."', font_size=13, italic=True, color=DARK)
add_textbox(slide, 0.5, 7.0, 12, 0.4, "KNUST × GridCo  ·  DLinear+H10  ·  PRODUCTION READY", font_size=10, bold=True, color=GREEN_LIGHT)

# Save
out_path = r"C:\Users\User\Desktop\project\llm_from_scratch-main\LOADFORECASINGPROJECT\dlinear_presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Size: {os.path.getsize(out_path) / 1024:.0f} KB")
print(f"Slides: {len(prs.slides)}")
